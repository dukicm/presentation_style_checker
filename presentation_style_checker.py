import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
import re

# Define your style rules here
STYLE_RULES = {
    "forbidden_patterns": [
        r"z\.\s*B\.",  # matches z.B., z. B., z.\u00A0B. etc.
        r"\bdu\b",
        "Du"
    ],
    "min_font_size": 11,
    "max_font_size": 22,
    "allowed_fonts":["Arial", "Verdana", "Frutiger"]
}

# Compile patterns once
COMPILED_PATTERNS = [re.compile(p) for p in STYLE_RULES["forbidden_patterns"]]

def get_all_shapes(shapes):
    all_shapes = []
    for shape in shapes:
        all_shapes.append(shape)
        if shape.shape_type == 6:  # 6 = GROUP
            all_shapes.extend(get_all_shapes(shape.shapes))
    return all_shapes

def check_slide(slide, slide_number, seen_results):
    results = []

    for shape in get_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()

            # Step 1: Build full paragraph text
            paragraph_text = ''.join(run.text for run in paragraph.runs)

            # Step 2: Search forbidden patterns in full paragraph
            for pattern in COMPILED_PATTERNS:
                matches = pattern.findall(paragraph_text)
                for match in matches:
                    for run in paragraph.runs:
                        if any(char in match for char in run.text):
                            run.font.color.rgb = RGBColor(255, 0, 0)
                    results.append([slide_number, "Forbidden Spelling", match, paragraph_text])

            # Check font size
            min_size = STYLE_RULES["min_font_size"]
            max_size = STYLE_RULES["max_font_size"]

            too_small = False
            too_large = False
            for run in paragraph.runs:
                font_size = run.font.size
                print(f"Slide {slide_number} → Text: {run.text} | Size: {font_size}")
                if font_size:
                    size_pt = font_size.pt
                    if size_pt < min_size:
                        run.font.color.rgb = RGBColor(255, 0, 0)
                        too_small = True
                    elif size_pt > max_size:
                        run.font.color.rgb = RGBColor(255, 0, 0)
                        too_large = True

            if too_small:
                results.append([slide_number, "Font Too Small", f"< {min_size}pt", paragraph.text])
            if too_large:
                results.append([slide_number, "Font Too Large", f"> {max_size}pt", paragraph.text])

    return results  # ✅ Correctly indented, only once, after all shapes are checked
    


def main():
    st.title("PowerPoint Style Checker")

    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

    if uploaded_file:
        prs = Presentation(uploaded_file)

        all_results = []
        seen_results = set()
        for i, slide in enumerate(prs.slides, start=1):
            all_results.extend(check_slide(slide, i, seen_results))

        # Save modified version
        output_filename = "marked_output.pptx"
        prs.save(output_filename)

        if all_results:
            st.subheader("Style Issues Found:")
            for result in all_results:
                st.write(f"Slide {result[0]}: ❌ {result[1]} → “{result[2]}” in: “{result[3]}”")
        else:
            st.success("✅ No style issues found!")

        # ✅ Download button for modified presentation
        with open(output_filename, "rb") as f:
            st.download_button(
                label="📥 Download Corrected Presentation",
                data=f,
                file_name=output_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

if __name__ == "__main__":
    main()

